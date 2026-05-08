from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Use default PowerPoint template with proper slide layouts
prs = Presentation()

# Professional Color Scheme
NAVY_BLUE = RGBColor(0, 31, 63)
DARK_GRAY = RGBColor(31, 41, 55)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    # Subtitle
    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.text = subtitle
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    return slide

def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    # Content
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = ""
    for bullet in bullets:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.level = 0

    return slide

def add_two_content_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two content layout

    # Title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    # Left content
    left_placeholder = slide.placeholders[1]
    left_placeholder.text = left_title
    left_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    left_placeholder.text_frame.paragraphs[0].font.bold = True
    left_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    for bullet in left_bullets:
        p = left_placeholder.text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0

    # Right content
    right_placeholder = slide.placeholders[2]
    right_placeholder.text = right_title
    right_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    right_placeholder.text_frame.paragraphs[0].font.bold = True
    right_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    for bullet in right_bullets:
        p = right_placeholder.text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0

    return slide

def add_diagram_slide(title, diagram_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

    # Diagram box
    diagram_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    diagram_frame = diagram_box.text_frame
    diagram_frame.text = diagram_text
    diagram_frame.paragraphs[0].font.size = Pt(11)
    diagram_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    diagram_frame.paragraphs[0].font.name = 'Courier New'

    return slide

# Slide 1: Title Slide
add_title_slide(
    "AI/ML Phishing Detection System",
    "Real-Time URL Analysis & Automated Alert System\nInternship Project Presentation"
)

# Slide 2: Agenda
add_bullet_slide(
    "Agenda",
    [
        "Introduction & Problem Statement",
        "Project Objectives & Technologies",
        "System Architecture & Data Flow",
        "Key Features & Implementation",
        "Results & Future Scope",
        "Conclusion"
    ]
)

# Slide 3: Introduction
add_bullet_slide(
    "Introduction",
    [
        "AI/ML-powered phishing URL detection system",
        "Complete web application with real-time analysis",
        "Dual alert system: Email + Browser notifications",
        "Interactive dashboard with threat visualization",
        "Built using modern web technologies and machine learning",
        "Internship project demonstrating practical AI implementation"
    ]
)

# Slide 4: Problem Statement
add_bullet_slide(
    "Problem Statement",
    [
        "Phishing attacks remain one of the most prevalent cyber threats",
        "Fake websites that mimic legitimate ones to steal credentials",
        "Manual detection is slow and unreliable for average users",
        "Current security solutions often fail against sophisticated attacks",
        "Need for automated, real-time detection using AI/ML",
        "Users require instant alerts when encountering suspicious URLs"
    ]
)

# Slide 5: Project Objectives
add_bullet_slide(
    "Project Objectives",
    [
        "Develop ML model for automatic phishing URL detection",
        "Create user-friendly web interface for URL analysis",
        "Implement real-time email and browser notifications",
        "Provide explainable AI with detailed threat analysis",
        "Build interactive threat map and performance dashboard",
        "Ensure end-to-end functionality from input to alert delivery"
    ]
)

# Slide 6: Technologies Used
add_two_content_slide(
    "Technologies Used",
    "Frontend Technologies",
    [
        "React 18 - Modern UI framework",
        "TypeScript - Type-safe development",
        "Vite - Fast build tool",
        "Axios - HTTP client",
        "Leaflet - Interactive maps",
        "CSS3 - Responsive styling"
    ],
    "Backend Technologies",
    [
        "Python 3.10 - Core language",
        "Flask - REST API framework",
        "Scikit-learn - ML library",
        "Gmail SMTP - Email service",
        "SQLite - Data storage",
        "NumPy/Pandas - Data processing"
    ]
)

# Slide 7: System Architecture
arch_diagram = """
┌─────────────────────────────────────────────────────────┐
│                    USER BROWSER                          │
│  ┌─────────────────────────────────────────────────┐     │
│  │ 7 Frontend Pages:                             │     │
│  │ • Dashboard (Overview & Metrics)              │     │
│  │ • AI/ML Detection (URL Analysis)              │     │
│  │ • Settings (Email & Notifications)            │     │
│  │ • ML Model Analytics (Performance)            │     │
│  │ • Feature Analysis (URL Components)           │     │
│  │ • Dataset Statistics (Training Data)          │     │
│  │ • Email Alert History                         │     │
│  └─────────────────────────────────────────────────┘     │
└─────────────────────┬───────────────────────────────────┘
                      │ HTTP Requests (Axios)
                      │ localhost:5175 → localhost:5000
                      ↓
┌─────────────────────────────────────────────────────────┐
│                 BACKEND SERVER (Flask)                  │
│  ┌─────────────────────────────────────────────────┐     │
│  │ 3 Main API Endpoints:                         │     │
│  │ • /api/analyze-url (ML Prediction)            │     │
│  │ • /api/ml-metrics (Model Performance)         │     │
│  │ • /api/dataset-stats (Training Data)          │     │
│  └─────────────────────────────────────────────────┘     │
└─────────────────────┼───────────────────────────────────┘
          ┌───────────┼───────────┬───────────┬───────────┐
          │           │           │           │           │
          ↓           ↓           ↓           ↓           ↓
┌─────────┐ ┌─────────┐ ┌─────────┐ ┌─────────┐ ┌─────────┐
│ ML Model│ │  Gmail  │ │ Browser │ │Database│ │Feature  │
│Logistic │ │  SMTP   │ │Notif    │ │Storage │ │Analysis │
│Regression│ │ Service │ │ API     │ │(Local) │ │Engine   │
│15 URLs  │ │(Emails) │ │(Pop-ups)│ │        │ │         │
└─────────┘ └─────────┘ └─────────┘ └─────────┘ └─────────┘
"""
add_diagram_slide("System Architecture", arch_diagram)

# Slide 8: Data Flow Diagram
data_flow = """
DATA FLOW DIAGRAM - URL Analysis Process

┌─────────────────┐
│   USER INPUT    │ ← User pastes URL
│ "example.com"   │
└─────────┬───────┘
          │
          ↓
┌─────────────────┐
│  FRONTEND       │ ← Validates URL format
│  VALIDATION     │
└─────────┬───────┘
          │
          ↓
┌─────────────────┐     ┌─────────────────┐
│  API REQUEST    │────▶│  BACKEND        │
│  /api/analyze   │     │  FLASK SERVER   │
│  (URL + Creds)  │     └─────────┬───────┘
└─────────────────┘               │
                                  ↓
┌─────────────────┐     ┌─────────────────┐
│  FEATURE        │────▶│  ML MODEL       │
│  EXTRACTION     │     │  PREDICTION     │
│  (Domain, Len,  │     │  (Logistic Reg) │
│   Subdomains)   │     └─────────┬───────┘
└─────────────────┘               │
                                  ↓
┌─────────────────┐     ┌─────────────────┐
│  PREDICTION     │────▶│  ALERTS SENT    │
│  RESULT         │     │  PARALLEL:      │
│  PHISHING 92%   │     │  • Email        │
└─────────┬───────┘     │  • Browser Pop-up│
          │             └─────────────────┘
          ↓
┌─────────────────┐
│  FRONTEND       │ ← Displays results
│  DISPLAY        │ ← Shows confidence %
│  RESULTS        │ ← Lists reasons
└─────────────────┘
"""
add_diagram_slide("Data Flow Diagram", data_flow)

# Slide 9: System Flowchart
flowchart = """
COMPLETE SYSTEM FLOWCHART

START
  │
  ↓
┌─────────────────┐     YES → ┌─────────────────┐
│ FIRST TIME?     │──────────▶│ SETUP PROCESS   │
│ (Settings done?)│           │ • Enter Gmail   │
└─────────┬───────┘           │ • App password  │
          │ NO                │ • Enable alerts │
          ↓                   └─────────────────┘
┌─────────────────┐
│ AI/ML DETECTION │ ← User goes to analysis page
│ PAGE            │
└─────────┬───────┘
          │
          ↓
┌─────────────────┐
│ ENTER URL       │ ← User pastes suspicious URL
│ "phishingsite.com"│
└─────────┬───────┘
          │
          ↓
┌─────────────────┐
│ CLICK ANALYZE   │     ┌─────────────────┐
│ BUTTON          │────▶│ BACKEND PROCESS │
│                 │     │ • Extract features│
└─────────────────┘     │ • Run ML model   │
                        │ • Calculate prob │
                        └─────────┬───────┘
                                  │
                                  ↓
                        ┌─────────────────┐
                        │ PREDICTION      │
                        │ RESULT          │
                        └─────────┬───────┘
                                 / \
                                /   \
                               /     \
                    PHISHING        SAFE
                      │              │
                      ↓              ↓
            ┌─────────────────┐    ┌─────────────────┐
            │ SEND ALERTS     │    │ SHOW SAFE      │
            │ • Email alert   │    │ MESSAGE        │
            │ • Browser popup │    └─────────────────┘
            │ • Log to history│
            └─────────────────┘
                      │
                      ↓
            ┌─────────────────┐
            │ DISPLAY RESULTS │
            │ • Red badge     │
            │ • Confidence %  │
            │ • Reasons list  │
            └─────────────────┘
"""
add_diagram_slide("System Flowchart", flowchart)

# Slide 10: Key Features
add_bullet_slide(
    "Key Features",
    [
        "Real-time URL analysis with instant results",
        "Dual alert system: Email + Browser notifications",
        "Explainable AI - Shows why URLs are suspicious",
        "Interactive threat map showing attack locations",
        "Comprehensive ML model performance dashboard",
        "Feature importance analysis for URL components",
        "Dataset statistics and training metrics",
        "Complete history of detected phishing attempts"
    ]
)

# Slide 11: How It Works
add_bullet_slide(
    "How It Works",
    [
        "Step 1: User configures Gmail credentials in Settings",
        "Step 2: User pastes suspicious URL in Detection page",
        "Step 3: Frontend sends URL to backend API for analysis",
        "Step 4: Backend extracts URL features (domain, length, etc.)",
        "Step 5: ML model analyzes features and predicts phishing probability",
        "Step 6: System sends dual alerts - email and browser notification",
        "Step 7: Results displayed with confidence percentage and reasons",
        "Step 8: Alert logged in history for future reference"
    ]
)

# Slide 12: Results & Performance
add_bullet_slide(
    "Results & Performance",
    [
        "ML Model Accuracy: 90% (9 out of 10 URLs correctly classified)",
        "Precision: 88% (88% of phishing alerts are legitimate threats)",
        "Recall: 92% (catches 92% of actual phishing URLs)",
        "Real-time analysis: Results delivered in 1-2 seconds",
        "Dual alerts: Both email and browser notifications working",
        "7 fully functional pages with professional UI",
        "Interactive threat map showing global attack patterns",
        "Complete end-to-end system working without errors"
    ]
)

# Slide 13: Future Enhancements
add_bullet_slide(
    "Future Enhancements",
    [
        "Train model on larger dataset (10,000+ URLs) for higher accuracy",
        "Implement additional ML algorithms (Random Forest, Neural Networks)",
        "Add database integration for persistent data storage",
        "Create admin panel for manual URL labeling and model retraining",
        "Develop public API for other applications to use the service",
        "Build mobile application for smartphone URL checking",
        "Add advanced features like URL screenshot analysis",
        "Implement user authentication and multi-user support"
    ]
)

# Slide 14: Conclusion
add_bullet_slide(
    "Conclusion",
    [
        "Successfully developed complete AI/ML phishing detection system",
        "Demonstrated practical application of machine learning in cybersecurity",
        "Achieved 90% accuracy with Logistic Regression on limited training data",
        "Implemented dual alert system ensuring instant user notifications",
        "Created professional web interface with comprehensive analytics",
        "System works reliably from URL input to alert delivery",
        "Project showcases skills in full-stack development and AI/ML",
        "Provides solid foundation for future enhancements and deployment"
    ]
)

# Slide 15: Project Summary
add_bullet_slide(
    "Project Summary",
    [
        "Development Time: 2 weeks (Internship project timeline)",
        "Codebase: ~2000+ lines across frontend and backend",
        "Technologies: React, TypeScript, Python, Flask, Scikit-learn",
        "Key Learning: End-to-end ML pipeline implementation",
        "Challenges: Gmail SMTP setup, browser notifications, ML metrics",
        "Testing: Manual testing with various phishing and legitimate URLs",
        "Deployment: Local environment (localhost:5175 frontend, localhost:5000 backend)",
        "Documentation: Complete README with setup and usage instructions"
    ]
)

# Slide 16: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[1])
title_placeholder = slide.shapes.title
title_placeholder.text = "Thank You!"
title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
title_placeholder.text_frame.paragraphs[0].font.bold = True
title_placeholder.text_frame.paragraphs[0].font.color.rgb = NAVY_BLUE

content_placeholder = slide.placeholders[1]
content_placeholder.text = "Questions & Discussion\n\nThis presentation demonstrates the practical application of AI/ML in cybersecurity through a complete phishing detection system."
content_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
content_placeholder.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

prs.save(r"c:\Phishing_AIML Project\Professional_Presentation.pptx")
print("✅ Professional PowerPoint presentation created successfully!")
print("✅ Uses proper PowerPoint slide layouts and templates")
print("✅ File: Professional_Presentation.pptx")
print("✅ Total slides: 16")
print("✅ Professional formatting with standard layouts")
print("✅ Complete with diagrams, flowcharts, and all content")
print("✅ Ready for internship viva presentation!")